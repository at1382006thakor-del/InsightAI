import os
from datetime import datetime
from sqlalchemy.orm import Session
from sqlalchemy import func
from typing import Tuple

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from ..database.models import Sale, Product, Customer

# Directory for reports
REPORTS_DIR = "reports"
SCRATCH_DIR = os.path.join(REPORTS_DIR, "scratch")
os.makedirs(SCRATCH_DIR, exist_ok=True)

def build_pptx_report(db: Session, report_type: str, markdown_summary: str) -> str:
    """Compiles statistics, drafts strategy slides and compiles PowerPoint report."""
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx package is required to generate PowerPoint reports.")

    prs = Presentation()
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    pptx_filename = f"InsightAI_{report_type.capitalize()}_Report_{timestamp}.pptx"
    pptx_path = os.path.join(REPORTS_DIR, pptx_filename)

    # 1. Slide 1: Cover Page
    slide_layout = prs.slide_layouts[0] # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "InsightAI Business Report"
    subtitle.text = f"Executive Strategy Slide Deck — {report_type.upper()}\nGenerated on {datetime.now().strftime('%B %d, %Y')}"

    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(9, 9, 11) # dark zinc-950

    # Clean text styles
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(139, 92, 246) # brand-500
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(161, 161, 170) # zinc-400

    # 2. Slide 2: Key Metrics Overview Table
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Key Performance Indicators Summary"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(139, 92, 246)
    
    # Set background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(9, 9, 11)

    # Query metrics
    total_rev = db.query(func.sum(Sale.revenue)).scalar() or 0.0
    total_prof = db.query(func.sum(Sale.profit)).scalar() or 0.0
    total_orders = db.query(func.count(Sale.sale_id)).scalar() or 0
    total_custs = db.query(func.count(Customer.customer_id)).scalar() or 0
    margin = (total_prof / total_rev * 100) if total_rev > 0 else 0.0

    # Add Table
    rows, cols = 5, 2
    left, top, width, height = Inches(1.5), Inches(2.0), Inches(7.0), Inches(3.5)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    metrics = [
        ("Total Sales Revenue", f"${total_rev:,.2f}"),
        ("Net Profit", f"${total_prof:,.2f}"),
        ("Profit Margin Percentage", f"{margin:.2f}%"),
        ("Closed Orders Volume", f"{total_orders:,}"),
        ("Active Client Base", f"{total_custs:,}")
    ]

    for idx, (metric_name, val) in enumerate(metrics):
        cell_name = table.cell(idx, 0)
        cell_val = table.cell(idx, 1)
        
        cell_name.text = metric_name
        cell_val.text = val
        
        # Style text inside cells
        for cell in (cell_name, cell_val):
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(244, 244, 245)
            p.font.name = "Arial"

    # 3. Slide 3: Recommendations Bullet List
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Strategic AI Recommendations"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(139, 92, 246)
    
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(9, 9, 11)

    tf = slide.placeholders[1].text_frame
    tf.text = "Based on transactional intelligence checks, we recommend the following:"

    # Parse strategy recommendations from Markdown summary
    bullet_points = []
    lines = markdown_summary.split("\n")
    for line in lines:
        line = line.strip()
        if (line.startswith("- ") or line.startswith("* ")) and len(line) > 15:
            clean_bullet = line[2:].replace("**", "").replace("*", "")
            bullet_points.append(clean_bullet)

    if bullet_points:
        for bp in bullet_points[:4]: # Max 4 main bullet points for layout
            p = tf.add_paragraph()
            p.text = bp
            p.level = 1
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(244, 244, 245)
    else:
        # Fallback suggestions
        fallbacks = [
            "Adjust regional pricing schemas in technology products to protect profit yields.",
            "Deactivate rule-based discount triggers on apparel lines showing margins below 5%.",
            "Establish automated low stock restock points to prevent inventory shortages."
        ]
        for f in fallbacks:
            p = tf.add_paragraph()
            p.text = f
            p.level = 1
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(244, 244, 245)

    prs.save(pptx_path)
    return pptx_path
